from pathlib import Path
import json

from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, ListFlowable, ListItem, Preformatted, Table, TableStyle
from pptx import Presentation

BASE = Path(r"c:\Users\JQY\Desktop\web1\cw1")
DOCS = BASE / "docs"
DOCS.mkdir(exist_ok=True)
openapi = json.loads((DOCS / "openapi.json").read_text(encoding="utf-8"))

styles = getSampleStyleSheet()
styles.add(ParagraphStyle(name="TitleCenter", parent=styles["Title"], alignment=TA_CENTER, fontSize=22, leading=26, spaceAfter=12))
styles.add(ParagraphStyle(name="HeadingSmall", parent=styles["Heading2"], fontSize=13, leading=16, spaceBefore=8, spaceAfter=6))
styles.add(ParagraphStyle(name="BodySmall", parent=styles["BodyText"], fontSize=9.5, leading=12))
styles.add(ParagraphStyle(name="CodeSmall", fontName="Courier", fontSize=8, leading=9))


def build_api_pdf():
    path = DOCS / "api_documentation.pdf"
    doc = SimpleDocTemplate(str(path), pagesize=landscape(A4), leftMargin=1.2 * cm, rightMargin=1.2 * cm, topMargin=1.2 * cm, bottomMargin=1.2 * cm)
    story = []
    story.append(Paragraph("API Documentation", styles["TitleCenter"]))
    story.append(Paragraph("Task Management API - generated from the project OpenAPI schema", styles["BodySmall"]))
    story.append(Spacer(1, 0.3 * cm))
    info_table = Table([
        ["Title", openapi.get("info", {}).get("title", "Task Management API")],
        ["Version", openapi.get("info", {}).get("version", "1.0.0")],
        ["Description", openapi.get("info", {}).get("description", "")],
    ], colWidths=[3.5 * cm, 22.5 * cm])
    info_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#E8EEF7")),
        ("BOX", (0, 0), (-1, -1), 0.5, colors.grey),
        ("INNERGRID", (0, 0), (-1, -1), 0.25, colors.grey),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
    ]))
    story.append(info_table)
    story.append(Spacer(1, 0.4 * cm))
    for endpoint, methods in openapi.get("paths", {}).items():
        for method, details in methods.items():
            story.append(Paragraph(f"{method.upper()} {endpoint}", styles["HeadingSmall"]))
            story.append(Paragraph(f"Summary: {details.get('summary', 'No summary provided')}", styles["BodySmall"]))
            parameters = details.get("parameters", [])
            if parameters:
                rows = [["Name", "In", "Required", "Type", "Description"]]
                for p in parameters:
                    rows.append([
                        p.get("name", ""), p.get("in", ""), str(p.get("required", False)), p.get("schema", {}).get("type", ""), p.get("description", "")
                    ])
                t = Table(rows, colWidths=[3 * cm, 2 * cm, 2.2 * cm, 2 * cm, 15 * cm])
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#D9EAF7")),
                    ("BOX", (0, 0), (-1, -1), 0.5, colors.grey),
                    ("INNERGRID", (0, 0), (-1, -1), 0.25, colors.grey),
                    ("FONTSIZE", (0, 0), (-1, -1), 8),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ]))
                story.append(t)
                story.append(Spacer(1, 0.2 * cm))
            if details.get("requestBody"):
                story.append(Paragraph("Request body schema:", styles["BodySmall"]))
                story.append(Preformatted(json.dumps(details["requestBody"], indent=2)[:2500], styles["CodeSmall"]))
            responses = details.get("responses", {})
            if responses:
                rr = [["Status", "Description"]] + [[code, resp.get("description", "")] for code, resp in responses.items()]
                rt = Table(rr, colWidths=[3 * cm, 21.7 * cm])
                rt.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E7F4E4")),
                    ("BOX", (0, 0), (-1, -1), 0.5, colors.grey),
                    ("INNERGRID", (0, 0), (-1, -1), 0.25, colors.grey),
                    ("FONTSIZE", (0, 0), (-1, -1), 8),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ]))
                story.append(rt)
                story.append(Spacer(1, 0.3 * cm))
    doc.build(story)


def build_report_pdf():
    path = DOCS / "Technical_Report_339919905.pdf"
    doc = SimpleDocTemplate(str(path), pagesize=A4, leftMargin=1.6 * cm, rightMargin=1.6 * cm, topMargin=1.5 * cm, bottomMargin=1.5 * cm)
    story = []
    def bullets(items, style_name="BodySmall"):
        return ListFlowable([ListItem(Paragraph(i, styles[style_name])) for i in items], bulletType="bullet", leftIndent=14)
    story.append(Paragraph("Task Management API Technical Report", styles["TitleCenter"]))
    story.append(Paragraph("Course Code: XJCO3011 | Student ID: 339919905 | Coursework 1 Submission", styles["BodySmall"]))
    story.append(Paragraph("GitHub Repository: https://github.com/339919905/web-api-project", styles["BodySmall"]))
    story.append(Spacer(1, 0.35 * cm))
    story.append(Paragraph("1. Introduction", styles["Heading1"]))
    story.append(Paragraph("This project implements a data-driven Web API for task management using Python and FastAPI. The goal was to build a maintainable backend system that satisfies the coursework requirements for CRUD functionality, database integration, structured JSON responses, appropriate HTTP status codes, and interactive API documentation. The selected data model is a task entity containing a title, description, completion state, due date, and timestamp metadata.", styles["BodySmall"]))
    story.append(Paragraph("2. Technology Selection and System Architecture", styles["Heading1"]))
    story.append(Paragraph("FastAPI was chosen because it enables rapid API development, automatic Swagger/OpenAPI generation, and strongly typed endpoint definitions. SQLAlchemy was selected as the ORM because it provides a structured and scalable way to define models and interact with relational data. SQLite was used for persistence because it is lightweight and easy to configure for coursework development. Pydantic was used to validate request bodies and ensure that API responses follow predictable schemas.", styles["BodySmall"]))
    story.append(Paragraph("The architecture follows a simple layered design: clients send HTTP requests to the FastAPI application, request data is validated through Pydantic schemas, route handlers execute application logic, SQLAlchemy manages persistence, and SQLite stores the final records. The codebase is organised into dedicated files for startup configuration, routing, schema definitions, database modelling, automated testing, and documentation assets.", styles["BodySmall"]))
    story.append(bullets([
        "Client -> FastAPI application -> SQLAlchemy ORM -> SQLite database",
        "main.py handles startup, middleware, and exception handlers",
        "routers/items.py contains task-related CRUD operations",
        "database.py defines the Task model and database session handling",
        "schemas.py defines request and response validation rules",
    ], "BodySmall"))
    story.append(Paragraph("3. Implementation of Core Functionality", styles["Heading1"]))
    story.append(Paragraph("The application delivers complete CRUD functionality for one data model, which is the central requirement of the coursework. Users can create tasks, retrieve all tasks, retrieve an individual task by identifier, fully update tasks with PUT, partially update selected fields with PATCH, and delete records. In addition to the minimum requirements, the API also includes pagination, title search, filtering by completion state, timestamp tracking, and consistent JSON error responses.", styles["BodySmall"]))
    story.append(bullets([
        "GET / returns a health-style response and documentation links",
        "GET /api/tasks returns a paginated list with page, limit, total, and items",
        "GET /api/tasks/{task_id} returns a single task or a 404 error",
        "POST /api/tasks creates a task and returns HTTP 201",
        "PUT /api/tasks/{task_id} fully replaces an existing task",
        "PATCH /api/tasks/{task_id} supports partial updates",
        "DELETE /api/tasks/{task_id} removes a task and returns a JSON confirmation message",
    ], "BodySmall"))
    story.append(Paragraph("4. Error Handling and Validation", styles["Heading1"]))
    story.append(Paragraph("The API includes custom exception handlers so that validation errors, HTTP errors, and unexpected failures return a consistent JSON structure. This improves the professionalism of the system because API consumers receive predictable responses instead of mixed error formats. Request validation is handled through typed Pydantic schemas, missing resources return 404 errors, and successful creation requests return 201 status codes.", styles["BodySmall"]))
    story.append(Paragraph("5. Testing Strategy", styles["Heading1"]))
    story.append(Paragraph("Automated testing was implemented using pytest together with FastAPI's TestClient. A dedicated temporary SQLite database is created for the tests so that the development database is not modified during test execution. This isolation makes the tests repeatable and reliable. The current test suite validates the root endpoint, CRUD operations, filtering, pagination, timestamp fields, 404 behaviour, and invalid request handling.", styles["BodySmall"]))
    story.append(Paragraph("6. Version Control and Development Process", styles["Heading1"]))
    story.append(Paragraph("The repository was maintained with Git and published to GitHub. The commit history is organised into logical stages that reflect the major development steps of the project, including repository setup, schema and database work, endpoint implementation, sample data support, testing, deployment support, and final submission preparation. This staged history provides evidence of a structured development workflow.", styles["BodySmall"]))
    story.append(bullets([
        "8ca73e4 Initialize coursework API repository",
        "a06c6a7 Add database configuration and task schemas",
        "a7cf785 Implement CRUD task API endpoints",
        "2e79950 Add sample seed data for local demonstrations",
        "10a11b3 Add automated API tests with isolated database",
        "4816579 Document project setup and deployment workflow",
    ], "BodySmall"))
    story.append(Paragraph("7. Reflection, Limitations, and Future Work", styles["Heading1"]))
    story.append(Paragraph("A key challenge in this coursework was balancing simplicity with code quality. The final solution aims to remain understandable for academic demonstration while still including features that improve realism, such as custom error handling, timestamps, Docker support, and automated tests. There are still limitations. SQLite is appropriate for local coursework use, but a production-ready deployment would benefit from PostgreSQL. The current API also has no authentication or user-specific permissions. Future improvements could include JWT authentication, sorting parameters, database migrations with Alembic, and continuous integration.", styles["BodySmall"]))
    story.append(Paragraph("8. Generative AI Use Statement", styles["Heading1"]))
    story.append(Paragraph("Generative AI tools were used during the development and documentation process. AI assistance was used to help scaffold the FastAPI project structure, draft CRUD-oriented code patterns, suggest testing strategies, support documentation drafting, and help organise the repository workflow for submission. However, the final system was reviewed, run locally, tested, adjusted for environment-specific issues, and prepared for submission through deliberate manual verification. The student remained responsible for understanding the implementation, checking correctness, selecting which changes to keep, and making final submission decisions.", styles["BodySmall"]))
    story.append(Paragraph("9. Conclusion", styles["Heading1"]))
    story.append(Paragraph("In summary, the Task Management API meets the key coursework requirements by providing a working database-backed Web API with full CRUD support, JSON responses, correct HTTP methods, automatic OpenAPI documentation, and evidence of testing and version control. The additional features included in the final project strengthen the submission and demonstrate an effort to move beyond the minimum implementation threshold.", styles["BodySmall"]))
    doc.build(story)


def build_presentation():
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Task Management API"
    title_slide.placeholders[1].text = "XJCO3011 Coursework 1\nStudent ID: 339919905\nFastAPI + SQLAlchemy + SQLite"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Project Aim and Scope"
    tf = slide.placeholders[1].text_frame
    tf.text = "Build a data-driven Web API that satisfies the coursework requirements."
    for text in [
        "One core resource: Task",
        "Full CRUD support with database persistence",
        "JSON responses, validation, and status codes",
        "Interactive documentation through /docs and /redoc",
    ]:
        p = tf.add_paragraph(); p.text = text; p.level = 1

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Architecture and Design"
    tf = slide.placeholders[1].text_frame
    tf.text = "The system follows a simple but modular backend architecture:"
    for text in [
        "Client -> FastAPI routes -> SQLAlchemy ORM -> SQLite database",
        "Pydantic validates request and response payloads",
        "main.py configures middleware, exception handling, and startup",
        "Routing, schemas, and database logic are separated into modules",
    ]:
        p = tf.add_paragraph(); p.text = text; p.level = 1

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Feature Set Delivered"
    tf = slide.placeholders[1].text_frame
    tf.text = "The API goes beyond the minimum CRUD requirement:"
    for text in [
        "Create, read, update, patch, and delete task records",
        "Pagination using page and limit query parameters",
        "Search and filter support on task listings",
        "Timestamp fields for auditability",
        "Standardized JSON error responses",
    ]:
        p = tf.add_paragraph(); p.text = text; p.level = 1

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Evidence of Version Control"
    tf = slide.placeholders[1].text_frame
    tf.text = "Git history shows a staged development process:"
    for text in [
        "8ca73e4 Initialize coursework API repository",
        "a06c6a7 Add database configuration and task schemas",
        "a7cf785 Implement CRUD task API endpoints",
        "2e79950 Add sample seed data for local demonstrations",
        "10a11b3 Add automated API tests with isolated database",
        "4816579 Document project setup and deployment workflow",
        "49b39ac Add submission report presentation and API documentation",
        "5b7417f Refine submission assets and clean repository history",
    ]:
        p = tf.add_paragraph(); p.text = text; p.level = 1

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Testing and Quality Assurance"
    tf = slide.placeholders[1].text_frame
    tf.text = "Testing was used as evidence that the API is reliable."
    for text in [
        "pytest and FastAPI TestClient were used",
        "10 automated tests cover core API behaviour",
        "A temporary SQLite database isolates tests from development data",
        "The suite covers CRUD, filtering, 404 responses, and validation errors",
    ]:
        p = tf.add_paragraph(); p.text = text; p.level = 1

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Reflection and Future Improvements"
    tf = slide.placeholders[1].text_frame
    tf.text = "The project meets the core coursework goals and provides a solid backend foundation."
    for text in [
        "Strengths: modular structure, testing, documentation, and clear API behaviour",
        "Limitation: SQLite is suitable for coursework but not ideal for higher-scale production use",
        "Missing feature: authentication and user ownership of tasks",
        "Future work: deploy to Render, add JWT auth, and migrate to PostgreSQL",
    ]:
        p = tf.add_paragraph(); p.text = text; p.level = 1
    prs.save(str(DOCS / "Presentation_XJCO3011.pptx"))


if __name__ == "__main__":
    build_api_pdf()
    build_report_pdf()
    build_presentation()
